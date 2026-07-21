"""Test dell'upload file nella knowledge base: parser (.pdf/.docx/.pptx/.xlsx)
e l'endpoint multipart POST /knowledge/ingest.
"""

from __future__ import annotations

import importlib
import io

import pytest

from etoro_bot.knowledge.ingest import MAX_UPLOAD_BYTES, ingest_upload
from etoro_bot.knowledge.parsers import UnsupportedFileTypeError, extract_text


# -- fixture di file minimi validi, generati in memoria ----------------------


def _make_docx(text: str) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(text)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # layout vuoto
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = text
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _make_xlsx(cell_a: str, cell_b: str) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.append([cell_a, cell_b])
    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def _make_pdf() -> bytes:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


# -- (a) extract_text per formato ---------------------------------------------


def test_extract_text_docx():
    content = _make_docx("Tesi di investimento su Apple.")
    text = extract_text("nota.docx", content)
    assert "Tesi di investimento su Apple." in text


def test_extract_text_pptx():
    content = _make_pptx("Slide su Microsoft buyback")
    text = extract_text("slides.pptx", content)
    assert "Slide su Microsoft buyback" in text
    assert text.startswith("Slide 1:")


def test_extract_text_xlsx():
    content = _make_xlsx("AAPL", "Buy")
    text = extract_text("dati.xlsx", content)
    assert "AAPL" in text
    assert "Buy" in text


def test_extract_text_pdf_no_exceptions():
    content = _make_pdf()
    text = extract_text("documento.pdf", content)
    assert isinstance(text, str)  # pagina bianca: nessuna eccezione, testo eventualmente vuoto


def test_extract_text_md_txt_still_supported():
    assert extract_text("note.md", b"# Titolo\n\ntesto") == "# Titolo\n\ntesto"
    assert extract_text("note.txt", b"testo semplice") == "testo semplice"


def test_extract_text_unsupported_extension():
    with pytest.raises(UnsupportedFileTypeError):
        extract_text("malware.exe", b"\x00\x01")


def test_extract_text_corrupted_docx_raises_clear_error():
    with pytest.raises(ValueError, match="impossibile leggere"):
        extract_text("rotto.docx", b"non e' davvero un docx")


# -- (b) limite dimensione -----------------------------------------------------


def test_ingest_upload_too_large():
    oversized = b"x" * (MAX_UPLOAD_BYTES + 1)
    with pytest.raises(ValueError, match="troppo grande"):
        ingest_upload("grande.txt", oversized)


# -- (c) endpoint FastAPI -------------------------------------------------------


@pytest.fixture()
def client(repo, pg_url, tmp_path, monkeypatch):
    monkeypatch.setenv("DATABASE_URL", pg_url)
    monkeypatch.setenv("KILL_SWITCH_DIR", str(tmp_path))
    monkeypatch.setenv("STATE_DIR", str(tmp_path))
    monkeypatch.setenv("KNOWLEDGE_BASE_DIR", str(tmp_path / "knowledge_base"))
    monkeypatch.setenv("DISABLE_SCHEDULER", "1")
    monkeypatch.delenv("ETORO_BOT_KILL", raising=False)

    from fastapi.testclient import TestClient

    from etoro_bot.api import server

    importlib.reload(server)  # ricostruisce app e cache col nuovo env
    server.get_repo.cache_clear()
    with TestClient(server.app) as tc:
        yield tc


def test_ingest_endpoint_docx(client, monkeypatch):
    from etoro_bot.knowledge.kb import KnowledgeBase

    # KnowledgeBase gira in modalità degradata (niente Qdrant nei test):
    # mockiamo add_news per verificare che il conteggio dei chunk passi correttamente.
    monkeypatch.setattr(KnowledgeBase, "add_news", lambda self, items: len(items))
    monkeypatch.setattr(KnowledgeBase, "ensure_collections", lambda self: None)

    content = _make_docx("Report trimestrale su NVIDIA.")
    resp = client.post(
        "/knowledge/ingest",
        files={
            "file": (
                "report.docx",
                content,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
        data={"tickers": "nvda, msft"},
    )
    assert resp.status_code == 200
    body = resp.json()
    assert body["filename"] == "report.docx"
    assert body["chunks_indexed"] == 1


def test_ingest_endpoint_rejects_unsupported_extension(client):
    resp = client.post(
        "/knowledge/ingest",
        files={"file": ("malware.exe", b"\x00\x01", "application/octet-stream")},
    )
    assert resp.status_code == 415


def test_ingest_endpoint_degraded_kb_returns_zero(client):
    # Nessun mock: KnowledgeBase è degradata (niente Qdrant) → add_news ritorna 0,
    # ma l'endpoint non deve crashare.
    content = _make_xlsx("Ticker", "Note")
    resp = client.post(
        "/knowledge/ingest",
        files={
            "file": (
                "dati.xlsx",
                content,
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert resp.status_code == 200
    assert resp.json()["chunks_indexed"] == 0
