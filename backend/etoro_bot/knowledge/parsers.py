"""Estrazione testo da file caricati per l'ingestione nella knowledge base.

Formati supportati: .pdf, .docx, .pptx, .xlsx (extra `rag`) e .md/.txt (sempre
disponibili, stdlib). Le librerie di parsing pesanti sono importate LAZY dentro
ogni funzione — stesso pattern di `kb.py` — così il pacchetto resta importabile
anche senza l'extra `rag` installato.
"""

from __future__ import annotations

from pathlib import PurePosixPath

SUPPORTED_SUFFIXES = (".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt")


class UnsupportedFileTypeError(ValueError):
    """Sollevata quando l'estensione del file non è tra quelle supportate."""

    def __init__(self, filename: str) -> None:
        self.filename = filename
        super().__init__(
            f"tipo di file non supportato per '{filename}': estensioni ammesse "
            f"{', '.join(SUPPORTED_SUFFIXES)}"
        )


def _suffix(filename: str) -> str:
    return PurePosixPath(filename).suffix.lower()


def extract_text(filename: str, content: bytes) -> str:
    """Estrae il testo da `content` in base all'estensione di `filename`.

    Solleva `UnsupportedFileTypeError` per estensioni non riconosciute; ogni
    parser cattura le eccezioni della libreria sottostante e le rilancia con
    un messaggio chiaro ("impossibile leggere <filename>: <motivo>").
    """
    suffix = _suffix(filename)
    if suffix == ".pdf":
        return _extract_pdf(filename, content)
    if suffix == ".docx":
        return _extract_docx(filename, content)
    if suffix == ".pptx":
        return _extract_pptx(filename, content)
    if suffix == ".xlsx":
        return _extract_xlsx(filename, content)
    if suffix in (".md", ".txt"):
        return _extract_plain_text(filename, content)
    raise UnsupportedFileTypeError(filename)


def _extract_plain_text(filename: str, content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError(f"impossibile leggere {filename}: {exc}") from exc


def _extract_pdf(filename: str, content: bytes) -> str:
    import io

    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)
    except Exception as exc:
        raise ValueError(f"impossibile leggere {filename}: {exc}") from exc


def _extract_docx(filename: str, content: bytes) -> str:
    import io

    try:
        import docx

        document = docx.Document(io.BytesIO(content))
        parts = [p.text for p in document.paragraphs if p.text]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text for cell in row.cells if cell.text]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)
    except Exception as exc:
        raise ValueError(f"impossibile leggere {filename}: {exc}") from exc


def _extract_pptx(filename: str, content: bytes) -> str:
    import io

    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text
            ]
            if texts:
                slides_text.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(slides_text)
    except Exception as exc:
        raise ValueError(f"impossibile leggere {filename}: {exc}") from exc


def _extract_xlsx(filename: str, content: bytes) -> str:
    import io

    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets_text = []
        try:
            for sheet in workbook.worksheets:
                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(v) for v in row if v is not None]
                    if cells:
                        rows_text.append(" | ".join(cells))
                if rows_text:
                    sheets_text.append(f"{sheet.title}:\n" + "\n".join(rows_text))
        finally:
            workbook.close()
        return "\n\n".join(sheets_text)
    except Exception as exc:
        raise ValueError(f"impossibile leggere {filename}: {exc}") from exc
